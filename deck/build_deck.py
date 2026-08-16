"""Generates deck/bookly-pitch-deck.pptx — the Part 2 solution pitch deck.

Run with: python deck/build_deck.py
Regenerate any time the prose below should change; this script is the
source of truth for the deck, not something hand-edited after export.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT_PATH = Path(__file__).parent / "bookly-pitch-deck.pptx"

# -- palette -----------------------------------------------------------------
BG = RGBColor(0xFA, 0xF6, 0xEF)
PANEL = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_TAN = RGBColor(0xEF, 0xE7, 0xDB)
INK = RGBColor(0x2B, 0x23, 0x20)
MUTED = RGBColor(0x7A, 0x6F, 0x66)
ACCENT = RGBColor(0x7A, 0x4A, 0x2B)
ACCENT_LIGHT = RGBColor(0xC7, 0x9A, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE3, 0xD9, 0xC9)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT = "Helvetica Neue"


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    bg._element.getparent().remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return slide


def add_text(slide, left, top, width, height, text, size=18, color=INK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return box


def add_kicker_and_title(slide, kicker, title, title_size=32):
    add_text(slide, Inches(0.6), Inches(0.35), Inches(8), Inches(0.35), kicker.upper(),
              size=13, color=ACCENT, bold=True)
    add_text(slide, Inches(0.6), Inches(0.68), Inches(11.5), Inches(0.8), title,
              size=title_size, color=INK, bold=True)


def add_rounded_box(slide, left, top, width, height, fill=PANEL, line_color=BORDER, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def add_page_number(slide, n):
    add_text(slide, SLIDE_W - Inches(0.9), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.3),
              str(n), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


# -- Slide 1: Thesis -----------------------------------------------------------

def build_thesis_slide(prs):
    slide = blank_slide(prs)
    add_kicker_and_title(slide, "Bookly Support Agent — Solution Pitch", "The Thesis", title_size=34)

    statement_box = add_rounded_box(slide, Inches(0.6), Inches(1.55), Inches(12.1), Inches(1.55),
                                     fill=ACCENT, line_color=ACCENT, radius=0.08)
    tf = statement_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "A support agent should never guess or fabricate when it can ask the customer or check a real system."
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT

    pillars = [
        (
            "Grounded in real data",
            "Every order, customer, and book fact comes from an MCP tool call against the live REST API — "
            "never from the model's memory. The system prompt explicitly forbids fabricating order numbers, "
            "statuses, tracking numbers, or policy details.",
        ),
        (
            "Ask, don't assume",
            "Returns require order, item (if multi-item), and reason before initiate_return ever fires. "
            "“My order” with more than one candidate gets a clarifying question, not a best guess.",
        ),
        (
            "Bounded scope, safely",
            "A dedicated guardrails layer screens every message in and out against configurable categories, "
            "redirecting anything outside Bookly support before it reaches a tool or the customer.",
        ),
    ]
    col_w = Inches(3.87)
    gap = Inches(0.24)
    left0 = Inches(0.6)
    top = Inches(3.35)
    height = Inches(3.35)
    for i, (title, body) in enumerate(pillars):
        left = left0 + i * (col_w + gap)
        card = add_rounded_box(slide, left, top, col_w, height, fill=PANEL, radius=0.06)
        add_text(slide, left + Inches(0.28), top + Inches(0.28), col_w - Inches(0.56), Inches(0.6),
                 title, size=17, color=ACCENT, bold=True)
        add_text(slide, left + Inches(0.28), top + Inches(0.95), col_w - Inches(0.56), height - Inches(1.2),
                 body, size=13.5, color=INK, line_spacing=1.15)

    add_page_number(slide, 1)
    return slide


# -- Slide 2: Architecture -----------------------------------------------------

def build_architecture_slide(prs):
    slide = blank_slide(prs)
    add_kicker_and_title(slide, "How It Works", "Architecture: One Request, End to End", title_size=30)

    # Flow diagram: 6 boxes left to right
    boxes = [
        ("Frontend", "Single-page chat UI\n(HTML/CSS/JS)"),
        ("Backend /\nOrchestrator", "Conversation loop,\nsession state, system prompt"),
        ("LLM\n(tool use)", "The model decides:\nanswer, ask, or call a tool"),
        ("MCP Server", "5 scoped tools —\nnever raw SQL/routes"),
        ("REST API", "Owns all validation;\nonly thing touching the DB"),
        ("Database", "SQLite: customers,\nbooks, orders"),
    ]
    n = len(boxes)
    top = Inches(1.75)
    box_h = Inches(1.35)
    margin = Inches(0.55)
    gap = Inches(0.22)
    total_w = SLIDE_W - 2 * margin
    box_w = Emu(int((total_w - gap * (n - 1)) / n))

    centers = []
    for i, (title, sub) in enumerate(boxes):
        left = margin + i * (box_w + gap)
        fill = ACCENT if i in (1, 2) else PANEL
        text_color = WHITE if i in (1, 2) else INK
        sub_color = RGBColor(0xE9, 0xD9, 0xC8) if i in (1, 2) else MUTED
        box = add_rounded_box(slide, left, top, box_w, box_h, fill=fill,
                               line_color=ACCENT if i in (1, 2) else BORDER, radius=0.1)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.16)
        tf.margin_right = Inches(0.16)
        tf.margin_top = Inches(0.14)
        lines = title.split("\n")
        for li, line in enumerate(lines):
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(14.5)
            run.font.bold = True
            run.font.color.rgb = text_color
            run.font.name = FONT
        for line in sub.split("\n"):
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(2)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(10)
            run.font.color.rgb = sub_color
            run.font.name = FONT
        centers.append((left + box_w // 2, left, left + box_w))

    # forward arrows between boxes — python-pptx connectors have no arrowhead
    # support, so use actual arrow autoshapes instead of lines.
    arrow_h = Inches(0.22)
    arrow_y = top + box_h // 2 - arrow_h // 2
    for i in range(n - 1):
        x0 = centers[i][2]
        x1 = centers[i + 1][1]
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x0, arrow_y, x1 - x0, arrow_h)
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT
        arrow.line.fill.background()
        arrow.shadow.inherit = False

    add_text(slide, margin, top + box_h + Inches(0.25), total_w, Inches(0.35),
             "←  Grounded reply flows back: DB result → API → MCP tool result → LLM's final answer → orchestrator → frontend",
             size=12.5, color=MUTED, align=PP_ALIGN.CENTER)

    # Four component callouts
    callouts = [
        ("Orchestration", "The conversation loop in backend/app/orchestrator.py: run the model, execute any "
                           "tool calls, feed results back, repeat until a final reply — capped at 5 round-trips."),
        ("Tools", "The MCP server (backend/app/mcp_client.py + mcp-server/) exposes 5 task-shaped tools "
                  "(search_books, get_customer, find_customer_orders, get_order_status, initiate_return), "
                  "not raw CRUD."),
        ("Memory", "Per-session conversation history, in-memory in the backend process. Nothing is persisted "
                   "across restarts — a prototype-scoped choice, called out for production in slide 4."),
        ("Prompts", "The system prompt (backend/app/prompts.py) scopes the agent to Bookly support, lists the "
                    "live tools, and hard-codes the no-fabrication / ask-don't-assume rules."),
    ]
    col_w = Inches(2.98)
    gap2 = Inches(0.2)
    left0 = Inches(0.55)
    ctop = Inches(4.55)
    cheight = Inches(2.35)
    for i, (title, body) in enumerate(callouts):
        left = left0 + i * (col_w + gap2)
        card = add_rounded_box(slide, left, ctop, col_w, cheight, fill=PANEL_TAN, radius=0.08)
        add_text(slide, left + Inches(0.2), ctop + Inches(0.2), col_w - Inches(0.4), Inches(0.4),
                 title, size=15, color=ACCENT, bold=True)
        add_text(slide, left + Inches(0.2), ctop + Inches(0.68), col_w - Inches(0.4), cheight - Inches(0.9),
                 body, size=11.5, color=INK, line_spacing=1.12)

    add_page_number(slide, 2)
    return slide


# -- Slide 3: Key decisions -----------------------------------------------------

def build_decisions_slide(prs):
    slide = blank_slide(prs)
    add_kicker_and_title(slide, "Why It's Built This Way", "Key Decisions & Trade-offs", title_size=30)

    decisions = [
        (
            "Enforce tool use via MCP\nrather than trusting model memory",
            "Chose",
            "Every order/customer/book fact must come from an MCP tool call; the system prompt forbids "
            "fabrication outright.",
            "Traded off",
            "Extra latency per turn (a real tool round-trip instead of an instant guess) and more services "
            "to keep online.",
            "Worth it because",
            "A hallucinated order number or refund status is a trust-destroying failure for a bookstore agent. "
            "Grounding is non-negotiable for anything the customer will act on.",
        ),
        (
            "A dedicated guardrails layer\n(separate classifier call)",
            "Chose",
            "A moderation-endpoint-style check — a separate, cheap OpenAI call — screens every inbound and "
            "outbound message against configurable categories.",
            "Traded off",
            "Added latency (two extra model calls per turn) and cost, versus no screening or a hand-rolled "
            "keyword filter.",
            "Worth it because",
            "Keyword filters are trivially bypassed and don't generalize; a config-driven classifier is tunable "
            "without a redeploy and fails closed on a parse error.",
        ),
        (
            "Layered DB / API / MCP /\norchestrator boundaries",
            "Chose",
            "Four boundaries, each layer touching only the one directly below it (DB only via API; API only "
            "via MCP; MCP only via the backend).",
            "Traded off",
            "Five containers and moving parts to build and deploy, versus one monolith, for a take-home-sized "
            "problem.",
            "Worth it because",
            "The LLM never sees SQL or raw routes — only scoped tools. That's also what lets the database move "
            "to managed Postgres later without touching the agent at all.",
        ),
    ]
    col_w = Inches(3.98)
    gap = Inches(0.15)
    left0 = Inches(0.55)
    top = Inches(1.55)
    height = Inches(5.6)
    label_color = ACCENT
    label_h = Inches(0.25)
    label_to_body_gap = Inches(0.05)
    body_h = Inches(0.85)
    item_to_item_gap = Inches(0.15)
    for i, d in enumerate(decisions):
        title, l1, b1, l2, b2, l3, b3 = d
        left = left0 + i * (col_w + gap)
        card = add_rounded_box(slide, left, top, col_w, height, fill=PANEL, radius=0.05)
        pad = Inches(0.25)
        add_text(slide, left + pad, top + Inches(0.22), col_w - 2 * pad, Inches(0.85),
                 title, size=15.5, color=INK, bold=True, line_spacing=1.05)

        y = top + Inches(1.15)
        for label, body in ((l1, b1), (l2, b2), (l3, b3)):
            add_text(slide, left + pad, y, col_w - 2 * pad, label_h, label.upper(),
                     size=10.5, color=label_color, bold=True)
            y += label_h + label_to_body_gap
            add_text(slide, left + pad, y, col_w - 2 * pad, body_h, body,
                     size=11, color=INK, line_spacing=1.1)
            y += body_h + item_to_item_gap

    add_page_number(slide, 3)
    return slide


# -- Slide 4: What I'd do differently -----------------------------------------

def build_next_steps_slide(prs):
    slide = blank_slide(prs)
    add_kicker_and_title(slide, "With More Time / A Production Context", "What I'd Do Differently", title_size=30)

    items = [
        ("Managed Postgres, not SQLite-on-a-volume",
         "Durable, concurrent-safe storage that survives a container restart — the compose file already "
         "isolates the DB behind the API layer specifically so this swap doesn't touch the agent."),
        ("Real customer identity verification",
         "Today, “what's your email” is the only check before disclosing order details. Production "
         "needs real auth before that trust boundary is crossed."),
        ("A human escalation path",
         "No way today to hand off to a person when the agent can't help or a customer is upset — a real "
         "deployment needs that safety valve."),
        ("RAG for policy/FAQ content",
         "Shipping/returns/password-reset text lives in config.yaml today. Fine for three policies; a real "
         "catalog of support content needs retrieval, not a growing system prompt."),
        ("Observability across the MCP boundary",
         "Right now it's log lines. Production needs traces spanning backend → MCP → API so a slow or "
         "wrong tool call is diagnosable."),
        ("A real eval set, not just scripted conversations",
         "The test suite proves the plumbing (tool wiring, guardrail fail-closed behavior) deterministically. "
         "Production needs eval coverage of the model's actual judgment across many phrasings."),
    ]
    col_w = Inches(6.0)
    row_h = Inches(1.62)
    gap_x = Inches(0.3)
    gap_y = Inches(0.18)
    left0 = Inches(0.55)
    top0 = Inches(1.65)
    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        left = left0 + col * (col_w + gap_x)
        top = top0 + row * (row_h + gap_y)
        card = add_rounded_box(slide, left, top, col_w, row_h, fill=PANEL_TAN, radius=0.08)
        num = add_rounded_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(0.42), Inches(0.42),
                               fill=ACCENT, line_color=ACCENT, radius=0.5)
        tf = num.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i + 1)
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = WHITE
        run.font.name = FONT
        add_text(slide, left + Inches(0.78), top + Inches(0.18), col_w - Inches(1.0), Inches(0.35),
                 title, size=14, color=INK, bold=True)
        add_text(slide, left + Inches(0.78), top + Inches(0.56), col_w - Inches(1.0), row_h - Inches(0.75),
                 body, size=11, color=MUTED, line_spacing=1.15)

    add_page_number(slide, 4)
    return slide


def main():
    prs = new_presentation()
    build_thesis_slide(prs)
    build_architecture_slide(prs)
    build_decisions_slide(prs)
    build_next_steps_slide(prs)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
